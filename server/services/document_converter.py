"""Document → Markdown conversion router.

This is the replacement for the old markitdown wrapper. We now route per
file type and per content shape:

    PDF   → PyMuPDF4LLM fast path. If it yields enough text, ship it.
            Otherwise fall back to Marker (Surya layout + OCR) for scans
            and complex layouts. Marker requires a one-time ~2GB model
            download — we surface that as a 409 to the API layer rather
            than blocking.
    DOCX  → mammoth (semantic HTML) → markdownify (HTML → markdown).
    PPTX  → python-pptx (custom slide-by-slide markdown emitter).
    MD    → passthrough.

Why this shape:

    * One tool can't do all of these well. markitdown was good but the
      PDF quality was poor on anything beyond plain text and it has no
      OCR — the upgrade we want.
    * Heavy deps (PyTorch, Surya) only get loaded when a real scanned
      PDF arrives, not on every DOCX import.
    * Each branch is independently swappable. If MinerU lands later for
      Chinese-heavy users, it slots in next to ``_pdf_with_marker``.
"""

from __future__ import annotations

import logging
import os
import tempfile
from contextlib import suppress
from pathlib import Path

from markdownify import markdownify

from services import marker_state

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Heuristic for "did the fast path produce enough text to call it done?"
# ---------------------------------------------------------------------------
# 40 chars/page is the threshold below which we assume it's a scan or an
# image-only PDF. Plain text PDFs typically yield thousands of chars/page;
# scans yield single digits (page numbers, stray OCR-poisoned glyphs).
# This is intentionally generous — we'd rather over-trust PyMuPDF4LLM and
# skip a 30s Marker run than the other way around.
_MIN_CHARS_PER_PAGE_FOR_FAST_PATH = 40


def is_pdf_native_text(path: Path) -> tuple[bool, int, int]:
    """Probe a PDF: is there enough embedded text to skip OCR?

    Returns (is_native, page_count, char_count) so the caller can log
    why the routing decision went the way it did.
    """
    import pymupdf  # PyMuPDF, brought in transitively by pymupdf4llm

    with pymupdf.open(path) as doc:
        page_count = len(doc)
        if page_count == 0:
            return False, 0, 0
        char_count = sum(len(page.get_text("text") or "") for page in doc)

    chars_per_page = char_count / page_count
    return chars_per_page >= _MIN_CHARS_PER_PAGE_FOR_FAST_PATH, page_count, char_count


# ---------------------------------------------------------------------------
# PDF — fast path
# ---------------------------------------------------------------------------


def _pdf_with_pymupdf4llm(path: Path) -> str:
    import pymupdf4llm

    return pymupdf4llm.to_markdown(str(path), show_progress=False) or ""


# ---------------------------------------------------------------------------
# PDF — fallback (Marker)
# ---------------------------------------------------------------------------

# Cached across requests. ``create_model_dict`` is expensive (loads several
# hundred MB of weights into memory) so we hold the artifact dict for the
# lifetime of the process once it's been built.
_marker_artifacts = None


def _get_marker_artifacts():
    global _marker_artifacts
    if _marker_artifacts is None:
        # Default to MPS on Apple Silicon — Marker reads TORCH_DEVICE on
        # init. Users can override (CUDA, CPU) via env if they need to.
        os.environ.setdefault("TORCH_DEVICE", "mps")
        from marker.models import create_model_dict

        _marker_artifacts = create_model_dict()
    return _marker_artifacts


def _pdf_with_marker(path: Path) -> str:
    from marker.converters.pdf import PdfConverter
    from marker.output import text_from_rendered

    converter = PdfConverter(artifact_dict=_get_marker_artifacts())
    rendered = converter(str(path))
    text, _, _ = text_from_rendered(rendered)
    return text or ""


async def convert_pdf(path: Path, *, force_ocr: bool = False) -> str:
    """Route a PDF to fast-path or Marker.

    By default we probe the PDF and pick the cheapest path that works.
    When ``force_ocr=True`` the caller has explicitly asked for the
    high-quality OCR pipeline (the "Import with OCR" menu item) — we
    skip the probe and go straight to Marker.

    Raises ``MarkerModelsRequiredError`` if Marker is needed but its
    weights haven't been downloaded yet.
    """
    import asyncio

    if force_ocr:
        logger.info("PDF routed to Marker (user requested OCR mode)")
        await marker_state.ensure_installed_or_409()
        return await asyncio.to_thread(_pdf_with_marker, path)

    is_native, pages, chars = await asyncio.to_thread(is_pdf_native_text, path)
    if is_native:
        logger.info(
            "PDF routed to PyMuPDF4LLM (fast): %d pages, %d chars (%.1f/page)",
            pages,
            chars,
            chars / max(pages, 1),
        )
        return await asyncio.to_thread(_pdf_with_pymupdf4llm, path)

    logger.info(
        "PDF routed to Marker (auto fallback — low text yield): %d pages, %d chars (%.1f/page)",
        pages,
        chars,
        chars / max(pages, 1),
    )
    await marker_state.ensure_installed_or_409()
    return await asyncio.to_thread(_pdf_with_marker, path)


# ---------------------------------------------------------------------------
# DOCX
# ---------------------------------------------------------------------------


def _docx_sync(path: Path) -> str:
    import mammoth

    with path.open("rb") as fh:
        result = mammoth.convert_to_html(fh)

    # Mammoth gives us tidy semantic HTML. markdownify turns it into
    # markdown that TipTap re-imports cleanly.
    md = markdownify(result.value, heading_style="ATX")
    return md or ""


async def convert_docx(path: Path) -> str:
    import asyncio

    return await asyncio.to_thread(_docx_sync, path)


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------


def _pptx_sync(path: Path) -> str:
    """Slide-by-slide markdown emitter.

    python-pptx exposes shapes per slide. We extract titles into H2,
    body placeholders into bullet lists or paragraphs, and tables into
    Markdown tables. Embedded images are referenced by name only — we
    don't re-upload them to ``~/.doxmind/uploads/`` here because the
    import flow doesn't yet have a hook for inline image extraction.
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(str(path))
    out: list[str] = []

    for idx, slide in enumerate(prs.slides, start=1):
        out.append(f"## Slide {idx}")

        title_text = None
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            title_text = (slide.shapes.title.text or "").strip()
            if title_text:
                out.append(f"### {title_text}")

        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue

            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if not line:
                        continue
                    indent = "  " * max(0, para.level)
                    out.append(f"{indent}- {line}")

            elif shape.has_table:
                rows = shape.table.rows
                if not rows:
                    continue
                headers = [
                    (cell.text or "").strip().replace("\n", " ") for cell in rows[0].cells
                ]
                out.append("| " + " | ".join(headers) + " |")
                out.append("|" + "|".join(["---"] * len(headers)) + "|")
                for row in list(rows)[1:]:
                    cells = [
                        (cell.text or "").strip().replace("\n", " ") for cell in row.cells
                    ]
                    out.append("| " + " | ".join(cells) + " |")

            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                name = shape.name or "image"
                out.append(f"_[{name}]_")

        out.append("")  # blank line between slides

    return "\n".join(out).strip()


async def convert_pptx(path: Path) -> str:
    import asyncio

    return await asyncio.to_thread(_pptx_sync, path)


# ---------------------------------------------------------------------------
# Router
# ---------------------------------------------------------------------------


async def convert_to_markdown(content: bytes, ext: str, *, force_ocr: bool = False) -> str:
    """Single entrypoint used by the import endpoint.

    ``ext`` includes the leading dot (".pdf", ".docx", ...).

    ``force_ocr`` is only meaningful for PDFs and asks the converter to
    skip the fast path and use Marker. The flag is silently ignored for
    other formats — there is no OCR to apply to a DOCX.
    """
    ext = ext.lower()

    with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
        tmp.write(content)
        tmp.flush()
        tmp_path = Path(tmp.name)

    try:
        if ext == ".pdf":
            return await convert_pdf(tmp_path, force_ocr=force_ocr)
        if ext == ".docx":
            return await convert_docx(tmp_path)
        if ext == ".pptx":
            return await convert_pptx(tmp_path)
        raise ValueError(f"Unsupported extension for converter router: {ext}")
    finally:
        with suppress(OSError):
            tmp_path.unlink()
